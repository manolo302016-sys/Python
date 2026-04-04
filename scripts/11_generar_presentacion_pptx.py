#!/usr/bin/env python3
"""
11_generar_presentacion_pptx.py — MentalPRO · Presentación Ejecutiva PowerPoint
=================================================================================
Genera una presentación ejecutiva (.pptx) de los resultados de riesgo psicosocial
para 1 empresa (Res. 2764/2022).

Uso:
    python scripts/11_generar_presentacion_pptx.py --empresa INDECOL

Salida:
    data/final/Presentacion_RiesgoPsicosocial_INDECOL_YYYYMMDD.pptx
"""

import sys
import io
import logging
import argparse
from pathlib import Path
from datetime import datetime

import pandas as pd
import numpy as np
import yaml

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Setup ─────────────────────────────────────────────────────────────────────
BASE_DIR = Path(__file__).resolve().parent.parent
log = logging.getLogger(__name__)

# ── Paleta R11 — inamovible ───────────────────────────────────────────────────
NAVY  = RGBColor(0x0A, 0x16, 0x28)
GOLD  = RGBColor(0xC9, 0x95, 0x2A)
CYAN  = RGBColor(0x00, 0xC2, 0xCB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY  = RGBColor(0xF3, 0xF4, 0xF6)
DARK  = RGBColor(0x1F, 0x29, 0x37)
MID   = RGBColor(0x6B, 0x72, 0x80)

RIESGO_RGB = {
    1: RGBColor(0x10, 0xB9, 0x81),
    2: RGBColor(0x6E, 0xE7, 0xB7),
    3: RGBColor(0xF5, 0x9E, 0x0B),
    4: RGBColor(0xF9, 0x73, 0x16),
    5: RGBColor(0xEF, 0x44, 0x44),
}

# matplotlib hex
MPL_RIESGO = {
    "Sin riesgo":      "#10B981",
    "Riesgo bajo":     "#6EE7B7",
    "Riesgo medio":    "#F59E0B",
    "Riesgo alto":     "#F97316",
    "Riesgo muy alto": "#EF4444",
}
MPL_NAVY = "#0A1628"
MPL_GOLD = "#C9952A"

ETIQUETAS = {
    1: "Sin riesgo", 2: "Riesgo bajo", 3: "Riesgo medio",
    4: "Riesgo alto", 5: "Riesgo muy alto",
}

# Diapositiva 16:9 estándar
W  = Inches(13.33)
H  = Inches(7.5)


# ── Helpers de color ──────────────────────────────────────────────────────────
def _rgb(r, g, b) -> RGBColor:
    return RGBColor(r, g, b)


def _semaforo_rgb(pct: float) -> RGBColor:
    if pct > 35:
        return RIESGO_RGB[5]
    if pct >= 15:
        return RIESGO_RGB[3]
    return RIESGO_RGB[1]


def _pct(v) -> str:
    if v is None or (isinstance(v, float) and np.isnan(v)):
        return "N/D"
    return f"{float(v):.1f}%"


# ── Helpers de figura matplotlib ─────────────────────────────────────────────
def _fig_to_buf(fig) -> io.BytesIO:
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


def _chart_stacked_h(df_src: pd.DataFrame, titulo: str = "") -> io.BytesIO:
    """Barras horizontales apiladas — distribución 5 niveles."""
    col_order = [c for c in ETIQUETAS.values() if c in df_src.columns]
    df = df_src[~df_src.get("confidencial", pd.Series(False)).fillna(False)].copy()
    if df.empty or not col_order:
        return None
    nombres = df["nombre_nivel"].tolist()
    n = len(nombres)
    fig, ax = plt.subplots(figsize=(9, max(2.5, n * 0.65 + 0.8)))
    left = np.zeros(n)
    for col in col_order:
        vals = df[col].fillna(0).values
        bars = ax.barh(nombres, vals, left=left, color=MPL_RIESGO[col], label=col, height=0.55)
        for bar, v in zip(bars, vals):
            if v >= 9:
                fc = "white" if col in ("Riesgo alto", "Riesgo muy alto") else MPL_NAVY
                ax.text(bar.get_x() + bar.get_width() / 2,
                        bar.get_y() + bar.get_height() / 2,
                        f"{v:.0f}%", ha="center", va="center",
                        fontsize=7, color=fc, fontweight="bold")
        left += vals
    ax.axvline(35, color="#EF4444", linestyle="--", lw=0.9, alpha=0.8, label="Umbral 35%")
    ax.set_xlim(0, 100)
    ax.set_xlabel("% de trabajadores", fontsize=9)
    ax.tick_params(axis="y", labelsize=8.5)
    ax.tick_params(axis="x", labelsize=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(loc="lower right", fontsize=7.5, framealpha=0.85)
    if titulo:
        ax.set_title(titulo, fontsize=10, fontweight="bold", color=MPL_NAVY, pad=7)
    fig.tight_layout(pad=0.6)
    return _fig_to_buf(fig)


def _chart_hbar(nombres: list, valores: list, titulo: str = "",
                umbral: float = 35.0) -> io.BytesIO:
    """Barras horizontales simples con semáforo."""
    if not nombres:
        return None
    n = len(nombres)
    fig, ax = plt.subplots(figsize=(9, max(2.5, n * 0.7 + 0.8)))
    colores = []
    for v in valores:
        if v > 50:
            colores.append(MPL_RIESGO["Riesgo muy alto"])
        elif v > 35:
            colores.append(MPL_RIESGO["Riesgo alto"])
        elif v >= 15:
            colores.append(MPL_RIESGO["Riesgo medio"])
        else:
            colores.append(MPL_RIESGO["Sin riesgo"])
    bars = ax.barh(nombres, valores, color=colores, height=0.55)
    for bar, v in zip(bars, valores):
        ax.text(bar.get_width() + 0.8, bar.get_y() + bar.get_height() / 2,
                f"{v:.1f}%", va="center", fontsize=8.5,
                color=MPL_NAVY, fontweight="bold")
    ax.axvline(umbral, color="#EF4444", linestyle="--",
               lw=0.9, alpha=0.8, label=f"Umbral {umbral:.0f}%")
    ax.set_xlim(0, max(valores) * 1.2 + 5)
    ax.set_xlabel("% Riesgo Alto + Muy Alto", fontsize=9)
    ax.tick_params(axis="y", labelsize=8.5)
    ax.tick_params(axis="x", labelsize=8)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(fontsize=7.5, framealpha=0.85)
    if titulo:
        ax.set_title(titulo, fontsize=10, fontweight="bold", color=MPL_NAVY, pad=7)
    fig.tight_layout(pad=0.6)
    return _fig_to_buf(fig)


def _chart_demo(categorias: list, valores: list, titulo: str = "") -> io.BytesIO:
    """Barras horizontales para demografía."""
    if not categorias:
        return None
    n = len(categorias)
    fig, ax = plt.subplots(figsize=(8, max(2.2, n * 0.55 + 0.6)))
    bars = ax.barh(categorias, valores, color=MPL_NAVY, height=0.5)
    for bar, v in zip(bars, valores):
        ax.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height() / 2,
                f"{v:.1f}%", va="center", fontsize=8.5, color=MPL_NAVY)
    ax.set_xlabel("% de trabajadores", fontsize=9)
    ax.set_xlim(0, max(valores) * 1.25 + 3)
    ax.tick_params(axis="y", labelsize=8.5)
    ax.tick_params(axis="x", labelsize=8)
    ax.spines[["top", "right"]].set_visible(False)
    if titulo:
        ax.set_title(titulo, fontsize=10, fontweight="bold", color=MPL_NAVY, pad=7)
    fig.tight_layout(pad=0.6)
    return _fig_to_buf(fig)


# ═══════════════════════════════════════════════════════════════════════════════
class PresentacionRiesgoPsicosocial:
    """Genera la presentación ejecutiva en PowerPoint."""

    def __init__(self, empresa: str):
        self.empresa = empresa.upper().strip()
        with open(BASE_DIR / "config" / "config.yaml", "r", encoding="utf-8") as f:
            self.cfg = yaml.safe_load(f)
        self.proc = BASE_DIR / self.cfg["paths"]["processed"]
        self.out_dir = BASE_DIR / "data" / "final"
        self.out_dir.mkdir(parents=True, exist_ok=True)
        self.prs = Presentation()
        self.prs.slide_width  = W
        self.prs.slide_height = H
        self._load_data()

    # ── Datos ─────────────────────────────────────────────────────────────────
    def _rp(self, fname: str) -> pd.DataFrame:
        p = self.proc / fname
        return pd.read_parquet(p) if p.exists() else pd.DataFrame()

    def _fe(self, df: pd.DataFrame) -> pd.DataFrame:
        if "empresa" in df.columns and not df.empty:
            return df[df["empresa"] == self.empresa].copy()
        return df

    def _load_data(self):
        log.info(f"Cargando datos para {self.empresa}...")
        self.baremo     = self._fe(self._rp("fact_scores_baremo.parquet"))
        self.benchmark  = self._fe(self._rp("fact_benchmark.parquet"))
        self.kpis       = self._fe(self._rp("fact_v3_kpis_globales.parquet"))
        self.demo       = self._fe(self._rp("fact_v3_demografia.parquet"))
        self.costos     = self._fe(self._rp("fact_v3_costos.parquet"))
        self.areas      = self._fe(self._rp("fact_v3_ranking_areas.parquet"))
        self.riesgo_emp = self._fe(self._rp("fact_riesgo_empresa.parquet"))
        self.trabajador = self._fe(self._rp("dim_trabajador.parquet"))
        self.n_eval = self.trabajador["cedula"].nunique() if not self.trabajador.empty else 0
        self.sector = (
            self.riesgo_emp["sector_rag"].iloc[0]
            if not self.riesgo_emp.empty else "[DATO NO DISPONIBLE]"
        )
        self.fecha = datetime.today().strftime("%B de %Y").capitalize()
        log.info(f"  Evaluados: {self.n_eval} | Sector: {self.sector}")

    # ── Distribución de riesgo ────────────────────────────────────────────────
    def _dist_riesgo(self, instrumentos: list, nivel_analisis: str) -> pd.DataFrame:
        df = self.baremo[
            (self.baremo["instrumento"].isin(instrumentos)) &
            (self.baremo["nivel_analisis"] == nivel_analisis)
        ].copy()
        if df.empty:
            return pd.DataFrame()
        grp = (
            df.groupby(["nombre_nivel", "nivel_riesgo", "etiqueta_nivel"])
            .agg(n=("cedula", "count")).reset_index()
        )
        tot = grp.groupby("nombre_nivel")["n"].sum().rename("n_total")
        grp = grp.merge(tot, on="nombre_nivel")
        grp["pct"] = (grp["n"] / grp["n_total"] * 100).round(1)
        piv = grp.pivot_table(
            index=["nombre_nivel", "n_total"],
            columns="etiqueta_nivel", values="pct",
            aggfunc="sum", fill_value=0,
        ).reset_index()
        piv.columns.name = None
        col_order = [c for c in ETIQUETAS.values() if c in piv.columns]
        piv = piv[["nombre_nivel", "n_total"] + col_order]
        piv["%A+MA"] = (piv.get("Riesgo alto", 0) + piv.get("Riesgo muy alto", 0)).round(1)
        piv["confidencial"] = piv["n_total"] < 5
        return piv.sort_values("%A+MA", ascending=False).reset_index(drop=True)

    # ── Helpers de diapositiva ────────────────────────────────────────────────
    def _blank_slide(self):
        layout = self.prs.slide_layouts[6]  # blank
        return self.prs.slides.add_slide(layout)

    def _fill_bg(self, slide, rgb: RGBColor):
        """Rellena el fondo de la diapositiva con un color sólido."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = rgb

    def _rect(self, slide, left, top, width, height, fill_rgb: RGBColor = None,
              line_rgb: RGBColor = None, line_width_pt: float = 0):
        """Añade un rectángulo de color."""
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            left, top, width, height,
        )
        if fill_rgb:
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
        else:
            shape.fill.background()
        if line_rgb:
            shape.line.color.rgb = line_rgb
            shape.line.width = Pt(line_width_pt)
        else:
            shape.line.fill.background()
        return shape

    def _txbox(self, slide, left, top, width, height,
               text: str = "", bold: bool = False, italic: bool = False,
               fontsize: float = 18, color: RGBColor = DARK,
               align=PP_ALIGN.LEFT, word_wrap: bool = True) -> "slide.shapes":
        """Añade un cuadro de texto."""
        txb = slide.shapes.add_textbox(left, top, width, height)
        tf = txb.text_frame
        tf.word_wrap = word_wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.bold = bold
        run.font.italic = italic
        run.font.size = Pt(fontsize)
        run.font.color.rgb = color
        return txb

    def _add_img(self, slide, buf: io.BytesIO,
                 left, top, width, height=None):
        """Embebe imagen desde BytesIO."""
        if buf is None:
            return
        if height:
            slide.shapes.add_picture(buf, left, top, width, height)
        else:
            slide.shapes.add_picture(buf, left, top, width)

    def _header_bar(self, slide, titulo: str, subtitulo: str = ""):
        """Barra superior NAVY con título."""
        self._rect(slide, 0, 0, W, Inches(1.1), fill_rgb=NAVY)
        self._txbox(slide, Inches(0.3), Inches(0.12), Inches(11), Inches(0.55),
                    text=titulo, bold=True, fontsize=22, color=WHITE, align=PP_ALIGN.LEFT)
        if subtitulo:
            self._txbox(slide, Inches(0.3), Inches(0.65), Inches(11), Inches(0.38),
                        text=subtitulo, fontsize=12, color=GOLD, align=PP_ALIGN.LEFT)

    def _footer(self, slide, num: int):
        """Pie de página con número y marca."""
        self._rect(slide, 0, H - Inches(0.32), W, Inches(0.32), fill_rgb=NAVY)
        self._txbox(slide, Inches(0.2), H - Inches(0.30), Inches(10), Inches(0.28),
                    text=f"Avantum Consulting  |  Informe Riesgo Psicosocial — {self.empresa}  |  Res. 2764/2022",
                    fontsize=7, color=GRAY, align=PP_ALIGN.LEFT)
        self._txbox(slide, W - Inches(0.8), H - Inches(0.30), Inches(0.6), Inches(0.28),
                    text=str(num), fontsize=7, color=GOLD, align=PP_ALIGN.RIGHT)

    def _kpi_box(self, slide, left, top, w, h,
                 valor: str, etiqueta: str, color: RGBColor = NAVY):
        """Cuadro KPI con valor grande y etiqueta pequeña."""
        self._rect(slide, left, top, w, h, fill_rgb=color)
        self._txbox(slide, left, top + Inches(0.12), w, h * 0.55,
                    text=valor, bold=True, fontsize=28, color=WHITE, align=PP_ALIGN.CENTER)
        self._txbox(slide, left, top + h * 0.60, w, h * 0.38,
                    text=etiqueta, fontsize=9, color=GRAY, align=PP_ALIGN.CENTER)

    # ══════════════════════════════════════════════════════════════════════════
    # DIAPOSITIVAS
    # ══════════════════════════════════════════════════════════════════════════

    def _slide_portada(self):
        """Diapositiva 1: Portada."""
        slide = self._blank_slide()
        self._fill_bg(slide, NAVY)

        # Banda GOLD horizontal
        self._rect(slide, 0, Inches(3.0), W, Inches(0.06), fill_rgb=GOLD)

        # Logo / nombre empresa
        self._txbox(slide, Inches(0.8), Inches(0.5), Inches(11.5), Inches(0.6),
                    text="AVANTUM CONSULTING", bold=True, fontsize=11, color=GOLD,
                    align=PP_ALIGN.CENTER)

        # Título principal
        self._txbox(slide, Inches(0.5), Inches(1.3), Inches(12.3), Inches(1.0),
                    text="Informe de Resultados", bold=True, fontsize=38,
                    color=WHITE, align=PP_ALIGN.CENTER)
        self._txbox(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(0.7),
                    text="Evaluacion de Factores de Riesgo Psicosocial",
                    fontsize=22, color=GOLD, align=PP_ALIGN.CENTER)

        # Empresa
        self._txbox(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.7),
                    text=self.empresa, bold=True, fontsize=30,
                    color=WHITE, align=PP_ALIGN.CENTER)

        # Metadata
        meta = f"Sector: {self.sector}   |   {self.n_eval} evaluados   |   {self.fecha}"
        self._txbox(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(0.4),
                    text=meta, fontsize=12, color=GRAY, align=PP_ALIGN.CENTER)

        # Normativa
        self._txbox(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.4),
                    text="Resolucion 2764 de 2022  |  Ministerio de Trabajo Colombia  |  Confidencial",
                    fontsize=9, color=MID, align=PP_ALIGN.CENTER)

    def _slide_agenda(self):
        """Diapositiva 2: Contenido."""
        slide = self._blank_slide()
        self._header_bar(slide, "Contenido de la Presentacion", "Estructura del informe")
        self._footer(slide, 2)

        items = [
            "01  Ficha Tecnica",
            "02  Perfil Demografico (ASIS)",
            "03  Resultados Intralaborales Globales",
            "04  Analisis por Dominios",
            "05  Dimensiones Criticas",
            "06  Extralaborales y Estres",
            "07  Mapa de Calor por Areas",
            "08  Ausentismo y Costo Economico",
            "09  Conclusiones",
            "10  Plan de Intervencion",
        ]
        col1 = items[:5]
        col2 = items[5:]
        for i, txt in enumerate(col1):
            self._txbox(slide, Inches(0.5), Inches(1.3 + i * 0.95), Inches(6),
                        Inches(0.75), text=txt, fontsize=13, color=DARK)
            self._rect(slide, Inches(0.3), Inches(1.3 + i * 0.95),
                       Inches(0.06), Inches(0.55), fill_rgb=GOLD)
        for i, txt in enumerate(col2):
            self._txbox(slide, Inches(7.0), Inches(1.3 + i * 0.95), Inches(6),
                        Inches(0.75), text=txt, fontsize=13, color=DARK)
            self._rect(slide, Inches(6.8), Inches(1.3 + i * 0.95),
                       Inches(0.06), Inches(0.55), fill_rgb=CYAN)

    def _slide_ficha_tecnica(self):
        """Diapositiva 3: Ficha técnica."""
        slide = self._blank_slide()
        self._header_bar(slide, "Ficha Tecnica", "Metodologia de evaluacion")
        self._footer(slide, 3)

        n_a = len(self.trabajador[self.trabajador["forma_intra"] == "A"]) if not self.trabajador.empty else 0
        n_b = len(self.trabajador[self.trabajador["forma_intra"] == "B"]) if not self.trabajador.empty else 0

        filas = [
            ("Empresa evaluada",     self.empresa),
            ("Sector economico",     self.sector),
            ("Total evaluados",      f"{self.n_eval} personas  ({n_a} Forma A  /  {n_b} Forma B)"),
            ("Instrumentos",         "Forma A, Forma B, Extralaboral, Cuestionario de Estres"),
            ("Baremos",              "Nacionales por tipo de cargo y sexo (2010 / 2022)"),
            ("Clasificacion",        "Sin riesgo | Bajo | Medio | Alto | Muy alto"),
            ("Marco normativo",      "Resolucion 2764/2022 y Resolucion 2646/2008"),
            ("Referente nacional",   "ENCST 2013-2021 (II y III Encuesta Nacional)"),
            ("Periodo evaluacion",   self.fecha),
        ]
        y0 = Inches(1.25)
        row_h = Inches(0.54)
        for i, (campo, valor) in enumerate(filas):
            bg = GRAY if i % 2 == 0 else WHITE
            self._rect(slide, Inches(0.3), y0 + i * row_h,
                       Inches(12.7), row_h - Inches(0.04), fill_rgb=bg)
            self._txbox(slide, Inches(0.4), y0 + i * row_h + Inches(0.05),
                        Inches(3.5), row_h, text=campo, bold=True, fontsize=10, color=NAVY)
            self._txbox(slide, Inches(3.9), y0 + i * row_h + Inches(0.05),
                        Inches(9.0), row_h, text=valor, fontsize=10, color=DARK)

    def _slide_kpis_globales(self):
        """Diapositiva 4: KPIs ejecutivos."""
        slide = self._blank_slide()
        self._header_bar(slide, "Resultados Globales — KPIs Ejecutivos",
                         "Resumen de exposicion al riesgo psicosocial")
        self._footer(slide, 4)

        # Extraer KPIs clave desde riesgo_emp y kpis
        kpis_data = []
        if not self.riesgo_emp.empty:
            for _, r in self.riesgo_emp.iterrows():
                nr = int(r.get("nivel_riesgo_empresa", 0))
                kpis_data.append((
                    ETIQUETAS.get(nr, str(nr)),
                    r["instrumento"],
                    RIESGO_RGB.get(nr, NAVY),
                ))

        if not self.kpis.empty:
            ama = (
                self.kpis[self.kpis["nivel_num"] >= 4]
                .groupby("kpi_grupo")["pct_alto_muy_alto"].max()
            )
            for kpi_grupo, pct in ama.items():
                kpis_data.append((_pct(pct), f"A+MA - {kpi_grupo}", _semaforo_rgb(pct)))

        # Mostrar hasta 6 KPIs en 2 filas
        kw = Inches(3.8)
        kh = Inches(1.8)
        gap = Inches(0.25)
        for i, (valor, etiqueta, color) in enumerate(kpis_data[:6]):
            col = i % 3
            row = i // 3
            x = Inches(0.3) + col * (kw + gap)
            y = Inches(1.35) + row * (kh + gap)
            self._kpi_box(slide, x, y, kw, kh, valor, etiqueta, color)

        # Nota reevaluación
        reevaluar = (
            self.riesgo_emp["debe_reevaluar_1año"].any()
            if not self.riesgo_emp.empty and "debe_reevaluar_1año" in self.riesgo_emp.columns
            else False
        )
        if reevaluar:
            self._rect(slide, Inches(0.3), Inches(6.8), Inches(12.7), Inches(0.42),
                       fill_rgb=RGBColor(0xFE, 0xF3, 0xC7))
            self._txbox(slide, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.36),
                        text="Alerta: La empresa debe realizar una reevaluacion en no mas de 12 meses (Res. 2764/2022)",
                        fontsize=9, color=RGBColor(0x92, 0x40, 0x0E), bold=True)

    def _slide_demo(self):
        """Diapositiva 5: Perfil demografico."""
        slide = self._blank_slide()
        self._header_bar(slide, "Perfil Demografico (ASIS)",
                         "Caracteristicas sociodemograficas de la poblacion evaluada")
        self._footer(slide, 5)

        df = self.demo
        if df.empty:
            self._txbox(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1),
                        text="[DATO NO DISPONIBLE]", fontsize=14, color=MID)
            return

        def _get_var(variable: str):
            sub = df[(df["variable"] == variable) &
                     (~df.get("confidencial", pd.Series(False)).fillna(False))]
            return sub["categoria"].tolist(), sub["pct"].fillna(0).tolist()

        # Gráfica 1: categoría de cargo
        cats1, vals1 = _get_var("categoria_cargo")
        if cats1:
            buf1 = _chart_demo(cats1, vals1, "Categoria de Cargo")
            if buf1:
                self._add_img(slide, buf1, Inches(0.3), Inches(1.2), Inches(6.2))

        # Gráfica 2: antigüedad
        cats2, vals2 = _get_var("antiguedad_empresa")
        if cats2:
            buf2 = _chart_demo(cats2, vals2, "Antiguedad en la Empresa")
            if buf2:
                self._add_img(slide, buf2, Inches(6.8), Inches(1.2), Inches(6.2))

        # Indicadores clave
        n_planta = self.n_eval
        if not self.costos.empty and "n_planta" in self.costos.columns:
            n_planta = int(self.costos["n_planta"].iloc[0])
        cobertura = round(self.n_eval / n_planta * 100, 1) if n_planta > 0 else 0

        self._txbox(slide, Inches(0.3), Inches(6.5), Inches(12.7), Inches(0.4),
                    text=f"Cobertura de evaluacion: {self.n_eval} de {n_planta} trabajadores ({cobertura}%)",
                    fontsize=10, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    def _slide_intra_dominios(self):
        """Diapositiva 6: Dominios intralaborales."""
        slide = self._blank_slide()
        self._header_bar(slide, "Resultados Intralaborales — Dominios",
                         "Distribucion de los 5 niveles de riesgo por dominio")
        self._footer(slide, 6)

        df_dom = self._dist_riesgo(["IntraA", "IntraB"], "dominio")
        if df_dom.empty:
            self._txbox(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1),
                        text="[DATO NO DISPONIBLE]", fontsize=14, color=MID)
            return

        buf = _chart_stacked_h(df_dom, "")
        if buf:
            self._add_img(slide, buf, Inches(0.3), Inches(1.15), Inches(9.5))

        # Tabla resumen lateral
        df_vis = df_dom[~df_dom["confidencial"]]
        y = Inches(1.2)
        self._rect(slide, Inches(9.9), y, Inches(3.1), Inches(0.38), fill_rgb=NAVY)
        self._txbox(slide, Inches(9.95), y + Inches(0.04), Inches(3.0), Inches(0.32),
                    text="Dominio   % A+MA", bold=True, fontsize=8.5, color=WHITE)
        for i, (_, r) in enumerate(df_vis.iterrows()):
            bg = GRAY if i % 2 == 0 else WHITE
            self._rect(slide, Inches(9.9), y + Inches(0.38) + i * Inches(0.44),
                       Inches(3.1), Inches(0.44), fill_rgb=bg)
            label = r["nombre_nivel"][:22] + "..." if len(r["nombre_nivel"]) > 22 else r["nombre_nivel"]
            self._txbox(slide, Inches(9.95),
                        y + Inches(0.38) + i * Inches(0.44) + Inches(0.05),
                        Inches(2.2), Inches(0.38), text=label, fontsize=8, color=DARK)
            sem_rgb = _semaforo_rgb(r["%A+MA"])
            self._rect(slide, Inches(12.0),
                       y + Inches(0.38) + i * Inches(0.44) + Inches(0.05),
                       Inches(0.9), Inches(0.32), fill_rgb=sem_rgb)
            self._txbox(slide, Inches(12.0),
                        y + Inches(0.38) + i * Inches(0.44) + Inches(0.05),
                        Inches(0.9), Inches(0.32),
                        text=_pct(r["%A+MA"]), bold=True, fontsize=8,
                        color=WHITE, align=PP_ALIGN.CENTER)

    def _slide_top_dimensiones(self):
        """Diapositiva 7: Top 5 dimensiones críticas."""
        slide = self._blank_slide()
        self._header_bar(slide, "Dimensiones Intralaborales Criticas",
                         "Top 5 dimensiones con mayor porcentaje de riesgo Alto + Muy Alto")
        self._footer(slide, 7)

        df_dim = self._dist_riesgo(["IntraA", "IntraB"], "dimension")
        if df_dim.empty:
            self._txbox(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1),
                        text="[DATO NO DISPONIBLE]", fontsize=14, color=MID)
            return

        top5 = df_dim[~df_dim["confidencial"]].head(5)
        buf = _chart_hbar(
            top5["nombre_nivel"].tolist(),
            top5["%A+MA"].tolist(),
            "% Riesgo Alto + Muy Alto por Dimension",
        )
        if buf:
            self._add_img(slide, buf, Inches(0.3), Inches(1.15), Inches(8.8))

        # Tarjetas de ranking
        y0 = Inches(1.2)
        card_h = Inches(1.0)
        for i, (_, r) in enumerate(top5.iterrows()):
            sem_rgb = _semaforo_rgb(r["%A+MA"])
            self._rect(slide, Inches(9.3), y0 + i * (card_h + Inches(0.08)),
                       Inches(3.7), card_h, fill_rgb=sem_rgb)
            self._txbox(slide, Inches(9.35), y0 + i * (card_h + Inches(0.08)) + Inches(0.04),
                        Inches(2.5), Inches(0.5),
                        text=r["nombre_nivel"], fontsize=8, color=WHITE, bold=True)
            self._txbox(slide, Inches(9.35), y0 + i * (card_h + Inches(0.08)) + Inches(0.48),
                        Inches(3.5), Inches(0.44),
                        text=_pct(r["%A+MA"]), fontsize=18, color=WHITE,
                        bold=True, align=PP_ALIGN.LEFT)

    def _slide_extralaboral(self):
        """Diapositiva 8: Extralaboral y Estrés."""
        slide = self._blank_slide()
        self._header_bar(slide, "Resultados Extralaborales y Estres",
                         "Distribucion de riesgo en factores externos al trabajo")
        self._footer(slide, 8)

        df_extra = self._dist_riesgo(["Extralaboral"], "dimension")
        df_estres = self._dist_riesgo(["Estres"], "dimension")

        # Gráfica extralaboral (izquierda)
        if not df_extra.empty:
            buf_e = _chart_stacked_h(df_extra, "Dimensiones Extralaborales")
            if buf_e:
                self._add_img(slide, buf_e, Inches(0.3), Inches(1.15), Inches(6.4))

        # Gráfica estrés (derecha)
        if not df_estres.empty:
            buf_s = _chart_stacked_h(df_estres, "Sintomas de Estres")
            if buf_s:
                self._add_img(slide, buf_s, Inches(6.9), Inches(1.15), Inches(6.1))

        if df_extra.empty and df_estres.empty:
            self._txbox(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1),
                        text="[DATO NO DISPONIBLE] — No se encontraron datos de extralaboral/estres.",
                        fontsize=12, color=MID)

    def _slide_mapa_calor(self):
        """Diapositiva 9: Mapa de calor por áreas."""
        slide = self._blank_slide()
        self._header_bar(slide, "Mapa de Calor por Areas",
                         "Porcentaje de trabajadores en Riesgo Alto + Muy Alto por area")
        self._footer(slide, 9)

        df = self.areas
        if df.empty:
            self._txbox(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1),
                        text="[DATO NO DISPONIBLE]", fontsize=14, color=MID)
            return

        df_vis = df[~df.get("confidencial", pd.Series(False)).fillna(False)].sort_values(
            "pct_ama", ascending=False
        )

        # Gráfica
        buf = _chart_hbar(
            df_vis["area_departamento"].tolist(),
            df_vis["pct_ama"].tolist(),
            "% Riesgo Alto + Muy Alto por Area",
            umbral=40.0,
        )
        if buf:
            self._add_img(slide, buf, Inches(0.3), Inches(1.15), Inches(8.2))

        # Tabla derecha
        y = Inches(1.15)
        self._rect(slide, Inches(8.6), y, Inches(4.4), Inches(0.42), fill_rgb=NAVY)
        self._txbox(slide, Inches(8.65), y + Inches(0.05), Inches(4.3), Inches(0.35),
                    text="Area   |   % A+MA   |   Nivel", bold=True,
                    fontsize=8, color=WHITE)
        for i, (_, r) in enumerate(df_vis.iterrows()):
            bg = GRAY if i % 2 == 0 else WHITE
            self._rect(slide, Inches(8.6), y + Inches(0.42) + i * Inches(0.5),
                       Inches(4.4), Inches(0.5), fill_rgb=bg)
            self._txbox(slide, Inches(8.65), y + Inches(0.42) + i * Inches(0.5) + Inches(0.05),
                        Inches(2.2), Inches(0.42),
                        text=r["area_departamento"], fontsize=8, color=DARK)
            sem_rgb = _semaforo_rgb(r["pct_ama"])
            self._rect(slide, Inches(10.9), y + Inches(0.42) + i * Inches(0.5) + Inches(0.06),
                       Inches(0.9), Inches(0.36), fill_rgb=sem_rgb)
            self._txbox(slide, Inches(10.9), y + Inches(0.42) + i * Inches(0.5) + Inches(0.06),
                        Inches(0.9), Inches(0.36),
                        text=_pct(r["pct_ama"]), bold=True, fontsize=8,
                        color=WHITE, align=PP_ALIGN.CENTER)

    def _slide_ausentismo(self):
        """Diapositiva 10: Ausentismo."""
        slide = self._blank_slide()
        self._header_bar(slide, "Ausentismo y Costo Economico",
                         "Indicadores de ausentismo laboral y estimacion de costo (Metodologia R10)")
        self._footer(slide, 10)

        df = self.costos
        if df.empty:
            self._txbox(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1),
                        text="[DATO NO DISPONIBLE]", fontsize=14, color=MID)
            return

        pct_aus = df["pct_ausentismo"].iloc[0] if "pct_ausentismo" in df.columns else 0
        dif = df["diferencia_pp_vs_pais"].iloc[0] if "diferencia_pp_vs_pais" in df.columns else 0

        # KPIs ausentismo
        sem = _semaforo_rgb(pct_aus) if pct_aus > 0 else RIESGO_RGB[1]
        self._kpi_box(slide, Inches(0.4), Inches(1.3), Inches(3.8), Inches(1.8),
                      _pct(pct_aus), "Tasa de ausentismo", sem)
        dif_rgb = RIESGO_RGB[5] if dif > 0 else RIESGO_RGB[1]
        self._kpi_box(slide, Inches(4.5), Inches(1.3), Inches(3.8), Inches(1.8),
                      f"{dif:+.1f}pp", "Diferencia vs pais", dif_rgb)
        n_planta = int(df["n_planta"].iloc[0]) if "n_planta" in df.columns else self.n_eval
        self._kpi_box(slide, Inches(8.6), Inches(1.3), Inches(3.8), Inches(1.8),
                      str(n_planta), "Trabajadores en planta", NAVY)

        # Tabla pasos
        y = Inches(3.3)
        self._rect(slide, Inches(0.4), y, Inches(12.5), Inches(0.4), fill_rgb=NAVY)
        self._txbox(slide, Inches(0.5), y + Inches(0.05), Inches(7), Inches(0.33),
                    text="Indicador", bold=True, fontsize=9, color=WHITE)
        self._txbox(slide, Inches(7.5), y + Inches(0.05), Inches(4), Inches(0.33),
                    text="Valor", bold=True, fontsize=9, color=WHITE, align=PP_ALIGN.CENTER)

        for i, (_, r) in enumerate(df.sort_values("paso").iterrows()):
            bg = GRAY if i % 2 == 0 else WHITE
            rh = Inches(0.42)
            self._rect(slide, Inches(0.4), y + Inches(0.4) + i * rh,
                       Inches(12.5), rh, fill_rgb=bg)
            self._txbox(slide, Inches(0.5), y + Inches(0.4) + i * rh + Inches(0.04),
                        Inches(7.0), rh, text=str(r.get("nombre_paso", "")),
                        fontsize=8.5, color=DARK)
            self._txbox(slide, Inches(7.5), y + Inches(0.4) + i * rh + Inches(0.04),
                        Inches(4.0), rh, text=str(r.get("valor", "")),
                        fontsize=8.5, color=DARK, align=PP_ALIGN.CENTER)

    def _slide_conclusiones(self):
        """Diapositiva 11: Conclusiones."""
        slide = self._blank_slide()
        self._header_bar(slide, "Conclusiones", "Hallazgos principales y recomendaciones")
        self._footer(slide, 11)

        # Extraer datos clave
        df_dim = self._dist_riesgo(["IntraA", "IntraB"], "dimension")
        top3 = df_dim[~df_dim.get("confidencial", pd.Series(False)).fillna(False)].head(3) \
            if not df_dim.empty else pd.DataFrame()

        riesgo_label = "[DATO NO DISPONIBLE]"
        if not self.riesgo_emp.empty:
            nr = int(self.riesgo_emp["nivel_riesgo_empresa"].iloc[0])
            riesgo_label = f"{ETIQUETAS.get(nr, str(nr))} (nivel {nr}/5)"

        puntos = [
            f"Nivel de riesgo global: {riesgo_label} para la poblacion evaluada ({self.n_eval} trabajadores).",
            f"Sector economico: {self.sector}.",
        ]
        if not top3.empty:
            dims = ", ".join(top3["nombre_nivel"].tolist())
            pcts = ", ".join([_pct(v) for v in top3["%A+MA"].tolist()])
            puntos.append(f"Dimensiones criticas: {dims} ({pcts} respectivamente).")
        if not self.areas.empty:
            top_area = self.areas.sort_values("pct_ama", ascending=False).iloc[0]
            puntos.append(
                f"Area con mayor criticidad: {top_area['area_departamento']} "
                f"({_pct(top_area['pct_ama'])} riesgo Alto+Muy Alto)."
            )
        if not self.costos.empty and "pct_ausentismo" in self.costos.columns:
            pct_a = self.costos["pct_ausentismo"].iloc[0]
            dif_a = self.costos["diferencia_pp_vs_pais"].iloc[0] if "diferencia_pp_vs_pais" in self.costos.columns else 0
            puntos.append(
                f"Ausentismo: {_pct(pct_a)} ({dif_a:+.1f}pp vs promedio nacional)."
            )
        puntos.append(
            "Conforme a la Res. 2764/2022, se requiere implementar el plan de "
            "intervencion en los plazos establecidos y realizar seguimiento periodico."
        )

        for i, punto in enumerate(puntos):
            self._rect(slide, Inches(0.3), Inches(1.25) + i * Inches(0.92),
                       Inches(0.06), Inches(0.62), fill_rgb=GOLD)
            self._txbox(slide, Inches(0.55), Inches(1.25) + i * Inches(0.92),
                        Inches(12.5), Inches(0.8), text=punto, fontsize=11, color=DARK)

    def _slide_plan_intervencion(self):
        """Diapositiva 12: Plan de intervención."""
        slide = self._blank_slide()
        self._header_bar(slide, "Plan de Intervencion Prioritaria",
                         "Dimensiones con >=35% de riesgo Alto + Muy Alto")
        self._footer(slide, 12)

        ESTRATEGIAS = {
            "Demandas cuantitativas":        "Redistribucion de carga laboral",
            "Demandas emocionales":          "Talleres de regulacion emocional",
            "Demandas de la jornada de trabajo": "Revision de turnos y descansos",
            "Control sobre el trabajo":      "Programa autonomia y participacion",
            "Caracteristicas del liderazgo": "Formacion en liderazgo transformacional",
            "Consistencia del rol":          "Clarificacion de roles y funciones",
            "Claridad de rol":               "Actualizacion de perfiles de cargo",
            "Reconocimiento y compensacion": "Revision esquema de compensacion",
            "Relaciones sociales en el trabajo": "Programa convivencia laboral",
            "Balance entre la vida laboral y familiar": "Flexibilizacion horaria",
        }
        PLAZOS = lambda pct: "0–3 meses" if pct > 50 else ("3–6 meses" if pct > 35 else "6–12 meses")

        # Recopilar dimensiones críticas
        criticas = []
        for instrs in [["IntraA", "IntraB"], ["Extralaboral"], ["Estres"]]:
            df_s = self._dist_riesgo(instrs, "dimension")
            if not df_s.empty:
                for _, r in df_s[~df_s["confidencial"]].iterrows():
                    if r["%A+MA"] >= 35:
                        criticas.append((r["nombre_nivel"], r["%A+MA"]))
        criticas.sort(key=lambda x: x[1], reverse=True)

        if not criticas:
            self._txbox(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1),
                        text="No se identificaron dimensiones con riesgo >= 35%. Mantener programa de vigilancia.",
                        fontsize=12, color=MID)
            return

        # Encabezado tabla
        y = Inches(1.2)
        cols = [("Dimension",          Inches(4.5)),
                ("% A+MA",             Inches(1.3)),
                ("Estrategia",         Inches(4.5)),
                ("Plazo",              Inches(1.5)),
                ("Responsable",        Inches(1.5))]
        x = Inches(0.3)
        self._rect(slide, x, y, Inches(13.3), Inches(0.4), fill_rgb=NAVY)
        for col_name, col_w in cols:
            self._txbox(slide, x + Inches(0.05), y + Inches(0.05), col_w, Inches(0.33),
                        text=col_name, bold=True, fontsize=8.5, color=WHITE)
            x += col_w

        row_h = Inches(0.48)
        for i, (dim, pct) in enumerate(criticas[:10]):
            bg = RIESGO_RGB[5] if pct > 50 else RIESGO_RGB[4]
            est = ESTRATEGIAS.get(dim, "Intervencion especifica — definir con SST")
            plazo = PLAZOS(pct)
            resp = "SST / RRHH"

            self._rect(slide, Inches(0.3), y + Inches(0.4) + i * row_h,
                       Inches(13.3), row_h, fill_rgb=GRAY if i % 2 == 0 else WHITE)
            # Indicador de color izquierdo
            self._rect(slide, Inches(0.3), y + Inches(0.4) + i * row_h,
                       Inches(0.08), row_h, fill_rgb=bg)

            vals = [dim, _pct(pct), est, plazo, resp]
            widths_list = [Inches(4.42), Inches(1.22), Inches(4.42), Inches(1.42), Inches(1.42)]
            cx = Inches(0.38)
            for val, cw in zip(vals, widths_list):
                self._txbox(slide, cx, y + Inches(0.4) + i * row_h + Inches(0.05),
                            cw, row_h - Inches(0.05),
                            text=val, fontsize=8, color=DARK)
                cx += cw

    def _slide_cierre(self):
        """Diapositiva 13: Cierre."""
        slide = self._blank_slide()
        self._fill_bg(slide, NAVY)
        self._rect(slide, 0, Inches(3.0), W, Inches(0.06), fill_rgb=GOLD)

        self._txbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.6),
                    text="AVANTUM CONSULTING", bold=True, fontsize=11,
                    color=GOLD, align=PP_ALIGN.CENTER)
        self._txbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.0),
                    text="Gracias", bold=True, fontsize=48,
                    color=WHITE, align=PP_ALIGN.CENTER)
        self._txbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.6),
                    text="Informe de Riesgo Psicosocial — " + self.empresa,
                    fontsize=16, color=GOLD, align=PP_ALIGN.CENTER)
        self._txbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
                    text="Resolucion 2764 de 2022  |  Ministerio de Trabajo Colombia",
                    fontsize=11, color=GRAY, align=PP_ALIGN.CENTER)
        self._txbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
                    text="Para consultas sobre este informe, contacte al equipo de SST de Avantum.",
                    fontsize=10, color=MID, align=PP_ALIGN.CENTER)

    # ── Generador principal ───────────────────────────────────────────────────
    def generar(self) -> Path:
        hoy = datetime.today().strftime("%Y%m%d")
        out = self.out_dir / f"Presentacion_RiesgoPsicosocial_{self.empresa}_{hoy}.pptx"

        log.info(f"Generando presentacion PPTX para {self.empresa}...")
        for slide_fn in [
            self._slide_portada,
            self._slide_agenda,
            self._slide_ficha_tecnica,
            self._slide_kpis_globales,
            self._slide_demo,
            self._slide_intra_dominios,
            self._slide_top_dimensiones,
            self._slide_extralaboral,
            self._slide_mapa_calor,
            self._slide_ausentismo,
            self._slide_conclusiones,
            self._slide_plan_intervencion,
            self._slide_cierre,
        ]:
            slide_fn()
            log.info(f"  Diapositiva: {slide_fn.__name__}")

        self.prs.save(str(out))
        log.info(f"PPTX guardado: {out}")
        return out


# ── CLI ───────────────────────────────────────────────────────────────────────
def main():
    logging.basicConfig(
        level=logging.INFO,
        format="%(levelname)s  %(message)s",
        handlers=[logging.StreamHandler(sys.stdout)],
    )
    parser = argparse.ArgumentParser(
        description="Genera presentacion ejecutiva PowerPoint de riesgo psicosocial."
    )
    parser.add_argument("--empresa", required=True,
                        help="Nombre de la empresa (ej: INDECOL)")
    args = parser.parse_args()

    pres = PresentacionRiesgoPsicosocial(empresa=args.empresa)
    out = pres.generar()
    print(f"\n{'='*60}")
    print(f"  Presentacion generada:")
    print(f"  {out}")
    print(f"{'='*60}\n")


if __name__ == "__main__":
    main()
